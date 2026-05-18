#!/usr/bin/env python3
"""
build-pptx-from-template.py

Build a PPTX deck by cloning a source template (pixel-perfect: same masters,
layouts, footers, logos) and injecting planned slide content into the
appropriate layouts.

Input: a slide plan JSON (schema below) and the source template PPTX.

Slide plan JSON schema:
{
  "metadata": {
    "title": "...", "subtitle": "...", "author": "...",
    "audience": "...", "tone": "..."
  },
  "template": "/abs/path/template.pptx",
  "output_pptx": "/abs/path/output.pptx",
  "slides": [
    {
      "role": "title" | "section" | "content" | "two_content" |
              "comparison" | "picture" | "title_only" | "blank",
      "layout_index": optional explicit override (within master 0),
      "title": "...",
      "subtitle": "...",
      "body": ["bullet 1", "bullet 2", ...],
      "body_left": [...], "body_right": [...],   // for two_content/comparison
      "images": [{"path": "/abs/path/img.png", "placeholder": "picture"}],
      "tables": [{"headers": [...], "rows": [[...], ...]}],
      "charts": [{
        "type": "bar"|"column"|"line"|"pie",
        "title": "...",
        "categories": ["A","B","C"],
        "series": [{"name": "...", "values": [1,2,3]}]
      }],
      "notes": "Speaker notes..."
    }
  ]
}

Usage:
  python3 scripts/build-pptx-from-template.py <plan.json> [--out <out.pptx>]
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu, Inches, Pt


CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}


def pick_layout(prs: Presentation, role: str, override_idx: int | None):
    """Find the best layout matching `role` from master 0."""
    master = prs.slide_masters[0]
    layouts = list(master.slide_layouts)
    if override_idx is not None and 0 <= override_idx < len(layouts):
        return layouts[override_idx]

    name_keywords = {
        "title": ["title slide", "title"],
        "section": ["section"],
        "two_content": ["two content", "two-content", "two"],
        "comparison": ["comparison"],
        "picture": ["picture", "caption"],
        "title_only": ["title only"],
        "blank": ["blank"],
        "content": ["title and content", "content"],
    }
    keywords = name_keywords.get(role, ["content"])
    for layout in layouts:
        name = (layout.name or "").lower()
        for kw in keywords:
            if kw in name:
                return layout

    # Fallbacks by index for the common 11-layout default set
    fallback_idx = {
        "title": 0,
        "section": 2,
        "content": 1,
        "two_content": 3,
        "comparison": 4,
        "title_only": 5,
        "blank": 6,
        "picture": 8,
    }.get(role, 1)
    fallback_idx = min(fallback_idx, len(layouts) - 1)
    return layouts[fallback_idx]


def get_placeholder_by_role(slide, *role_keywords: str):
    """Return the first placeholder whose name or type contains any keyword."""
    keywords = [k.lower() for k in role_keywords]
    for ph in slide.placeholders:
        type_name = ""
        try:
            type_name = str(ph.placeholder_format.type).split(".")[-1].lower()
        except Exception:
            pass
        name = (ph.name or "").lower()
        for kw in keywords:
            if kw in type_name or kw in name:
                return ph
    return None


def get_title_placeholder(slide):
    return get_placeholder_by_role(slide, "title", "center_title", "ctrtitle")


def get_subtitle_placeholder(slide):
    return get_placeholder_by_role(slide, "subtitle", "subtitle 2", "subtitle 1")


def get_body_placeholders(slide):
    """All non-title, non-subtitle, non-footer placeholders that can hold body text."""
    bodies = []
    for ph in slide.placeholders:
        try:
            type_name = str(ph.placeholder_format.type).split(".")[-1].lower()
        except Exception:
            type_name = ""
        if any(t in type_name for t in ("title", "subtitle", "footer", "slide_number", "date")):
            continue
        if ph.has_text_frame:
            bodies.append(ph)
    return bodies


def set_text_frame(tf, lines: list[str]):
    """Replace a text frame's content with a list of lines (each becomes a bullet/paragraph)."""
    if not lines:
        return
    # Preserve first paragraph's formatting if possible
    tf.clear()
    first = tf.paragraphs[0]
    first.text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line


def add_image(prs, slide, image_path: str, left=None, top=None, width=None, height=None):
    p = Path(image_path)
    if not p.exists():
        print(f"  warning: image not found: {image_path}", file=sys.stderr)
        return None
    prs_width = prs.slide_width
    prs_height = prs.slide_height
    if left is None:
        left = int(prs_width * 0.55)
    if top is None:
        top = int(prs_height * 0.20)
    if width is None:
        width = int(prs_width * 0.40)
    return slide.shapes.add_picture(str(p), left, top, width=width)


def add_table(prs, slide, headers, rows, left=None, top=None, width=None, height=None):
    prs_width = prs.slide_width
    prs_height = prs.slide_height
    if left is None:
        left = int(prs_width * 0.10)
    if top is None:
        top = int(prs_height * 0.30)
    if width is None:
        width = int(prs_width * 0.80)
    if height is None:
        height = int(prs_height * 0.50)

    cols = max(len(headers), max((len(r) for r in rows), default=1))
    n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(n_rows, cols, left, top, width, height)
    table = table_shape.table
    for c, h in enumerate(headers):
        table.cell(0, c).text = str(h)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            if c < cols:
                table.cell(r, c).text = str(val)
    return table_shape


def add_chart(prs, slide, spec: dict):
    chart_type = CHART_TYPE_MAP.get(spec.get("type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    data = CategoryChartData()
    data.categories = spec.get("categories", [])
    for s in spec.get("series", []):
        data.add_series(s.get("name", "Series"), s.get("values", []))

    prs_width = prs.slide_width
    prs_height = prs.slide_height
    left = int(prs_width * 0.10)
    top = int(prs_height * 0.25)
    width = int(prs_width * 0.80)
    height = int(prs_height * 0.60)

    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, data)
    if spec.get("title"):
        chart = graphic_frame.chart
        chart.has_title = True
        chart.chart_title.text_frame.text = spec["title"]
    return graphic_frame


def set_notes(slide, notes_text: str):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def populate_slide(prs: Presentation, slide_def: dict):
    role = slide_def.get("role", "content")
    layout = pick_layout(prs, role, slide_def.get("layout_index"))
    slide = prs.slides.add_slide(layout)

    # Title
    title_ph = get_title_placeholder(slide)
    if title_ph and slide_def.get("title"):
        title_ph.text_frame.text = slide_def["title"]

    # Subtitle
    subtitle_ph = get_subtitle_placeholder(slide)
    if subtitle_ph and slide_def.get("subtitle"):
        subtitle_ph.text_frame.text = slide_def["subtitle"]

    # Body content
    bodies = get_body_placeholders(slide)

    if role in ("two_content", "comparison"):
        left_lines = slide_def.get("body_left") or []
        right_lines = slide_def.get("body_right") or []
        if bodies:
            set_text_frame(bodies[0].text_frame, left_lines)
        if len(bodies) > 1:
            set_text_frame(bodies[1].text_frame, right_lines)
        elif right_lines:
            # Layout has only one body placeholder; add a text box on the right half.
            from pptx.util import Inches
            left = int(prs.slide_width * 0.52)
            top = int(prs.slide_height * 0.20)
            width = int(prs.slide_width * 0.43)
            height = int(prs.slide_height * 0.65)
            tb = slide.shapes.add_textbox(left, top, width, height)
            set_text_frame(tb.text_frame, right_lines)
    else:
        body_lines = slide_def.get("body") or []
        if bodies and body_lines:
            set_text_frame(bodies[0].text_frame, body_lines)

    # Images
    for img in slide_def.get("images", []) or []:
        path = img.get("path")
        if not path:
            continue
        add_image(prs, slide, path)

    # Tables
    for tbl in slide_def.get("tables", []) or []:
        add_table(prs, slide, tbl.get("headers", []), tbl.get("rows", []))

    # Charts
    for chart in slide_def.get("charts", []) or []:
        add_chart(prs, slide, chart)

    # Speaker notes
    set_notes(slide, slide_def.get("notes", ""))

    return slide


def remove_existing_slides(prs: Presentation):
    """Remove all pre-existing slides from the cloned template, including their parts and rels."""
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slide_rId_pairs = [(s, s.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"))
                       for s in list(sldIdLst)]
    pres_part = prs.part
    for sld_id_el, rId in slide_rId_pairs:
        sldIdLst.remove(sld_id_el)
        if rId:
            try:
                pres_part.drop_rel(rId)
            except Exception:
                pass


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("plan", help="Path to slide plan JSON")
    parser.add_argument("--out", help="Override output path")
    parser.add_argument(
        "--keep-template-slides",
        action="store_true",
        help="Keep slides already in the template (default: remove them)",
    )
    args = parser.parse_args()

    plan_path = Path(args.plan).resolve()
    if not plan_path.exists():
        print(f"Plan not found: {plan_path}", file=sys.stderr)
        return 1

    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    template = Path(plan["template"]).resolve()
    if not template.exists():
        print(f"Template not found: {template}", file=sys.stderr)
        return 1

    output = Path(args.out or plan["output_pptx"]).resolve()
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(template))
    if not args.keep_template_slides:
        remove_existing_slides(prs)

    # Apply core metadata
    meta = plan.get("metadata", {}) or {}
    if meta.get("title"):
        prs.core_properties.title = meta["title"]
    if meta.get("author"):
        prs.core_properties.author = meta["author"]
    if meta.get("audience"):
        prs.core_properties.subject = f"Audience: {meta['audience']}"

    # Build slides
    for i, slide_def in enumerate(plan.get("slides", []), start=1):
        populate_slide(prs, slide_def)
        print(f"  built slide {i}: role={slide_def.get('role', 'content')!r} title={slide_def.get('title', '')!r}")

    prs.save(str(output))
    print(f"Wrote PPTX: {output}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
