#!/usr/bin/env python3
"""
extract-pptx-template.py

Deep extraction of a PowerPoint template into a JSON manifest:
- Theme colors and fonts
- Slide size
- Every slide master + every layout
- Placeholders (type, index, name, position, size, default text)
- Layout role mapping (title, section, content, two_content, comparison, picture, blank)
- Core/app metadata

Usage:
  python3 scripts/extract-pptx-template.py <input.pptx> <output-manifest.json>
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# python-pptx placeholder type constants are enums; serialize as str
PLACEHOLDER_ROLE_HINTS = {
    "TITLE": "title",
    "CENTER_TITLE": "title",
    "SUBTITLE": "subtitle",
    "BODY": "body",
    "OBJECT": "content",
    "CONTENT": "content",
    "PICTURE": "picture",
    "TABLE": "table",
    "CHART": "chart",
    "MEDIA_CLIP": "media",
    "DATE": "date",
    "FOOTER": "footer",
    "SLIDE_NUMBER": "slide_number",
    "HEADER": "header",
    "BITMAP": "picture",
}


def emu_to_px(emu: int | None) -> float | None:
    if emu is None:
        return None
    # 914400 EMU per inch; render at 96 dpi -> px
    return round(Emu(emu).inches * 96, 2)


def placeholder_info(ph) -> dict:
    fmt = ph.placeholder_format
    t = fmt.type if fmt else None
    # python-pptx exposes PP_PLACEHOLDER enum; .name gives the clean key (e.g. "BODY").
    # str(t) would yield "BODY (2)" and break the role_hint lookup.
    type_name = getattr(t, "name", None) if t is not None else None
    role_hint = PLACEHOLDER_ROLE_HINTS.get(type_name) if type_name else None
    return {
        "idx": fmt.idx if fmt else None,
        "type": type_name,
        "role_hint": role_hint,
        "name": ph.name,
        "left_emu": ph.left,
        "top_emu": ph.top,
        "width_emu": ph.width,
        "height_emu": ph.height,
        "left_px": emu_to_px(ph.left),
        "top_px": emu_to_px(ph.top),
        "width_px": emu_to_px(ph.width),
        "height_px": emu_to_px(ph.height),
        "has_text": ph.has_text_frame,
        "default_text": (ph.text_frame.text if ph.has_text_frame else "") if ph.has_text_frame else "",
    }


def classify_layout(name: str, placeholders: list[dict]) -> str:
    """Map a layout to a semantic role."""
    n = (name or "").lower()
    types = [p["type"] or "" for p in placeholders]
    role_hints = [p["role_hint"] or "" for p in placeholders]
    body_like = sum(1 for r in role_hints if r in ("body", "content"))

    if "title slide" in n or "CENTER_TITLE" in types:
        return "title"
    if "section" in n:
        return "section"
    if "comparison" in n:
        return "comparison"
    if "two" in n or body_like >= 2:
        return "two_content"
    if "picture" in n or "PICTURE" in types or "BITMAP" in types:
        return "picture"
    if "title only" in n:
        return "title_only"
    if "blank" in n:
        return "blank"
    if "caption" in n:
        return "caption"
    return "content"


def theme_from_master(master) -> dict:
    """Best-effort theme color/font extraction from the master's referenced theme part."""
    out: dict = {"colors": {}, "fonts": {}}
    try:
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        from lxml import etree

        root = etree.fromstring(theme_part.blob)
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        clr_scheme = root.find(f".//{ns_a}clrScheme")
        if clr_scheme is not None:
            for child in clr_scheme:
                tag = etree.QName(child).localname
                srgb = child.find(f"{ns_a}srgbClr")
                sysclr = child.find(f"{ns_a}sysClr")
                if srgb is not None:
                    out["colors"][tag] = "#" + srgb.get("val", "").upper()
                elif sysclr is not None:
                    out["colors"][tag] = "#" + (sysclr.get("lastClr") or "").upper()

        font_scheme = root.find(f".//{ns_a}fontScheme")
        if font_scheme is not None:
            major = font_scheme.find(f"{ns_a}majorFont/{ns_a}latin")
            minor = font_scheme.find(f"{ns_a}minorFont/{ns_a}latin")
            if major is not None:
                out["fonts"]["major"] = major.get("typeface")
            if minor is not None:
                out["fonts"]["minor"] = minor.get("typeface")
    except Exception:
        pass

    return out


def core_props(prs: Presentation) -> dict:
    cp = prs.core_properties
    return {
        "title": cp.title or "",
        "author": cp.author or "",
        "subject": cp.subject or "",
        "keywords": cp.keywords or "",
        "language": cp.language or "",
        "category": cp.category or "",
        "comments": cp.comments or "",
    }


def main() -> int:
    if len(sys.argv) < 3:
        print("Usage: extract-pptx-template.py <input.pptx> <output-manifest.json>", file=sys.stderr)
        return 1

    src = Path(sys.argv[1]).resolve()
    dst = Path(sys.argv[2]).resolve()

    if not src.exists():
        print(f"Not found: {src}", file=sys.stderr)
        return 1

    prs = Presentation(str(src))

    masters_out = []
    role_map: dict[str, list[dict]] = {}

    for m_idx, master in enumerate(prs.slide_masters):
        layouts = []
        for l_idx, layout in enumerate(master.slide_layouts):
            ph_list = [placeholder_info(ph) for ph in layout.placeholders]
            role = classify_layout(layout.name, ph_list)
            layout_entry = {
                "master_index": m_idx,
                "layout_index": l_idx,
                "name": layout.name,
                "role": role,
                "placeholders": ph_list,
            }
            layouts.append(layout_entry)
            role_map.setdefault(role, []).append(
                {"master_index": m_idx, "layout_index": l_idx, "name": layout.name}
            )

        masters_out.append(
            {
                "master_index": m_idx,
                "name": master.name if hasattr(master, "name") else f"Master {m_idx + 1}",
                "theme": theme_from_master(master),
                "layouts": layouts,
            }
        )

    manifest = {
        "source_pptx": str(src),
        "slide_size_emu": {"width": prs.slide_width, "height": prs.slide_height},
        "slide_size_px": {
            "width": emu_to_px(prs.slide_width),
            "height": emu_to_px(prs.slide_height),
        },
        "aspect": (
            "16:9"
            if abs((prs.slide_width / prs.slide_height) - (16 / 9)) < 0.01
            else "4:3"
            if abs((prs.slide_width / prs.slide_height) - (4 / 3)) < 0.01
            else "custom"
        ),
        "core_properties": core_props(prs),
        "slide_count": len(prs.slides),
        "masters": masters_out,
        "role_map": role_map,
    }

    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(json.dumps(manifest, indent=2, default=str), encoding="utf-8")

    print(f"Wrote template manifest: {dst}")
    print(f"  Masters: {len(masters_out)}")
    print(f"  Total layouts: {sum(len(m['layouts']) for m in masters_out)}")
    print(f"  Roles found: {sorted(role_map.keys())}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
